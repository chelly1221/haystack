from pptx import Presentation
import re
import hashlib


def extract_cleaned_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    cleaned_lines = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    cleaned_lines.extend(text.splitlines())

    return cleaned_lines


def split_pptx_by_token_window(pptx_path, window_size=700, overlap=100, tokenizer=None):
    lines = extract_cleaned_text_from_pptx(pptx_path)
    all_text = "\n".join(lines)

    tokens = tokenizer.encode(all_text, add_special_tokens=False)
    total_tokens = len(tokens)
    chunks = []

    start = 0
    idx = 1
    while start < total_tokens:
        end = min(start + window_size, total_tokens)
        token_chunk = tokens[start:end]
        text_chunk = tokenizer.decode(token_chunk, skip_special_tokens=True)

        chunks.append({
            "title": f"Chunk {idx}",
            "content": text_chunk,
            "start_token": start
        })

        idx += 1
        start += window_size - overlap

    return chunks


def split_pptx_by_section_headings(pptx_path, heading_pattern=None):
    def parse_number(n):
        return [int(p) for p in n.split(".") if p.isdigit()]

    def is_strictly_next(prev, current):
        if not prev:
            return current == [1, 1, 1]
        if len(current) == len(prev):
            return current[:-1] == prev[:-1] and current[-1] == prev[-1] + 1
        if len(current) == len(prev) + 1:
            return current[:-1] == prev and current[-1] == 1
        for i in range(min(len(prev), len(current))):
            if current[i] == prev[i]:
                continue
            if current[i] == prev[i] + 1:
                return all(c == 1 for c in current[i + 1:])
            return False
        return False

    def hash_text(text):
        return hashlib.sha256(text.strip().encode("utf-8")).hexdigest()

    lines = extract_cleaned_text_from_pptx(pptx_path)

    if heading_pattern is None:
        heading_pattern = re.compile(r'^(\d{1,2}(?:\.\d{1,2}){0,3})\s+.+')

    sections = []
    collect = False
    last_number = None
    seen = set()
    merged_hashes = set()

    current_title = None
    current_lines = []

    for line in lines:
        match = heading_pattern.match(line)
        if match:
            section_number = match.group(1)
            section_number_list = parse_number(section_number)

            if not collect:
                if section_number == "1.1.1":
                    collect = True
                    last_number = section_number_list
                    current_title = line
                    current_lines = [line]
                continue

            if not is_strictly_next(last_number, section_number_list):
                if section_number in seen:
                    continue
                content_hash = hash_text("\n".join(current_lines))
                if content_hash in merged_hashes:
                    continue
                if sections:
                    sections[-1]["content"] += "\n\n" + "\n".join(current_lines)
                    merged_hashes.add(content_hash)
                    seen.add(section_number)
                continue

            if current_title and current_lines:
                sections.append({
                    "title": current_title,
                    "section_id": current_title.split()[0],
                    "content": "\n".join(current_lines)
                })

            current_title = line
            current_lines = [line]
            last_number = section_number_list
            seen.add(section_number)

        else:
            current_lines.append(line)

    if current_title and current_lines:
        sections.append({
            "title": current_title,
            "section_id": current_title.split()[0],
            "content": "\n".join(current_lines)
        })

    return sections
